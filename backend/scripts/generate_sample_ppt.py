from pptx import Presentation
from pptx.util import Inches

import os


def add_slide(prs, title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12), Inches(0.9))
    box.text_frame.paragraphs[0].text = title
    box.text_frame.paragraphs[0].font.size = 320000
    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(12), Inches(5))
    for i, line in enumerate(lines):
        para = body.text_frame.paragraphs[0] if i == 0 else body.text_frame.add_paragraph()
        para.text = line
        para.font.size = 240000


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_slide(prs, "北京代表处 2026Q2 经营汇报", ["机关领导：", "本材料为北京代表处 2026Q2 经营汇报，请审阅。"])
    add_slide(prs, "经营概况", ["本期营业收入保持增长，客户拓展顺利，整体经营平稳。"])
    add_slide(prs, "收入情况", ["营业收入 1000万元", "上年同期 864万元", "收入同比 +8.6%"])
    add_slide(prs, "利润情况", ["净利润 120万元", "利润同比 +15.7%"])
    add_slide(prs, "现金流", ["经营活动现金流 -50万元", "受应收账款增加影响"])
    add_slide(prs, "费用情况", ["销售费用 200万元", "管理费用 80万元"])
    add_slide(prs, "库存情况", ["库存 300万元", "库存环比 +40%"])
    add_slide(prs, "风险提示", ["海外订单存在延期风险，客户回款周期变长。"])
    add_slide(prs, "下一步计划", ["渠道整改事项按计划推进，预计三季度完成。"])
    add_slide(prs, "经营总结", ["营业收入 1200万元", "本期经营目标达成情况良好。"])

    out_dir = r"D:\A_code\AI\fin_agent\backend\data\samples"
    os.makedirs(out_dir, exist_ok=True)
    out = os.path.join(out_dir, "北京代表处2026Q2经营汇报.pptx")
    prs.save(out)
    print(out)


if __name__ == "__main__":
    main()
